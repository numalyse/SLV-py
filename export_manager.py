from PySide6.QtWidgets import QWidget, QVBoxLayout, QPushButton, QFileDialog, QRadioButton, QLabel, QLineEdit, QDialog, QButtonGroup, QHBoxLayout, QApplication
from PySide6.QtCore import Qt

import json
import os
import cv2
import numpy as np
import tempfile
import textwrap
import io

import moviepy.config as mp_config 
from moviepy import VideoFileClip, AudioFileClip
from moviepy.video.io.VideoFileClip import VideoFileClip
from moviepy.audio.io.AudioFileClip import AudioFileClip
from proglog import ProgressBarLogger
from PIL import Image, ImageDraw, ImageFont
from io import BytesIO
import imageio_ffmpeg 

# Pour PDF
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import Paragraph, SimpleDocTemplate
from reportlab.platypus import Image as PDFImage
# Pour PPTX
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor as PPTXColor
# Pour DOCX
from docx import Document  
from docx.shared import Pt, Inches, RGBColor as DOCXColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
# Pour ODT
from odf.opendocument import OpenDocumentText  
from odf.text import P, LineBreak
from odf.style import Style, TextProperties, ParagraphProperties, GraphicProperties
from odf.draw import Frame
from odf.draw import Image as ODTImage

# Importation de classes
from message_popup import MessagePopUp
from time_manager import TimeManager
from exportvideo_thread import ExportVideoThread
from exporttext_thread import ExportTextThread
from exporttagimages_thread import ExportTagImagesThread
from no_focus_push_button import NoFocusPushButton

import logging
import shutil, time

ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
os.environ["IMAGEIO_FFMPEG_EXE"] = ffmpeg_path

class ExportManager(QWidget):
    # Définition des styles
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'TitleStyle',
        parent=styles['Title'],
        fontSize=18,
        textColor=colors.red,
        spaceAfter=20
    )
    subtitle_style = ParagraphStyle(
        'SubtitleStyle',
        parent=styles['Heading2'],
        fontSize=12,
        textColor=colors.blue,
        spaceAfter=10
    )
    note_style = ParagraphStyle(
        'NoteStyle',
        parent=styles['BodyText'],
        fontSize=10,
        textColor=colors.black,
        leftIndent=20
    )

    def __init__(self, parent=None, vlc=None, projectmanager=None, format_export_text=[]):
        super().__init__(parent)
        self.seg = parent
        self.vlc = vlc
        self.project_manager = projectmanager

        self.file_path = self.project_manager.project_path
        self.title = self.project_manager.project_name
        self.time_manager = TimeManager(fps=self.vlc.fps)
        self.format_export_text = format_export_text

        self.configure()

    def configure(self):
        """ Ouvre une fenêtre de configuration pour choisir le mode. """
        dialog = QDialog(self)
        dialog.setWindowTitle("Exportation du Travail")
        dialog_layout = QVBoxLayout(dialog)

        # Options pour le type d'exportation
        num_label = QLabel("Type d'exportation :", dialog)
        dialog_layout.addWidget(num_label)

        num_group = QButtonGroup(dialog)
        option_1 = QRadioButton("Fichier texte", dialog)
        option_2 = QRadioButton("Vidéo annotée", dialog)
        option_3 = QRadioButton("TagImages", dialog)
        num_group.addButton(option_1)
        num_group.addButton(option_2)
        num_group.addButton(option_3)
        option_1.setChecked(True)
        dialog_layout.addWidget(option_1)
        dialog_layout.addWidget(option_2)
        dialog_layout.addWidget(option_3)

        dialog_load = QHBoxLayout()
        load = QLabel("")
        load.setStyleSheet("color: blue;")
        load.setAlignment(Qt.AlignCenter)
        dialog_load.addWidget(load)
        dialog_layout.addLayout(dialog_load)

        # Boutons OK/Annuler
        button_layout = QHBoxLayout()
        ok_button = NoFocusPushButton("OK", dialog)
        cancel_button = NoFocusPushButton("Annuler", dialog)
        button_layout.addWidget(ok_button)
        button_layout.addWidget(cancel_button)
        dialog_layout.addLayout(button_layout)

        def on_ok():
            if self.file_path:
                ok_button.setEnabled(False)
                load.setText("Exportation en cours ⌛")
                QApplication.processEvents()
                if option_1.isChecked():
                    if self.format_export_text[0]:
                        self.exported_thread = ExportTextThread(self,option=1)
                    elif self.format_export_text[1]:
                        self.exported_thread = ExportTextThread(self,option=2)
                    elif self.format_export_text[2]:
                        self.exported_thread = ExportTextThread(self,option=3)
                    elif self.format_export_text[3]:
                        self.exported_thread = ExportTextThread(self,option=4)
                elif option_2.isChecked():
                    self.default_file_path = os.path.join(self.file_path, f"{self.title}.mp4")
                    self.chosen_file_path, _ = QFileDialog.getSaveFileName(self, "Enregistrer l'exportation sous...", self.file_path, "Vidéo MP4 (*.mp4)")
                    self.exported_thread = ExportVideoThread(self)
                else:
                    self.exported_thread = ExportTagImagesThread(self)

                self.exported_thread.segmentation_done.connect(lambda: export_done(dialog))
                self.exported_thread.start()

            else:
                ok_button.setEnabled(True)

        def on_cancel():
            if hasattr(self, 'exported_thread'):
                print('Annulation')
                self.exported_thread.stop()
            dialog.reject()

        def export_done(dialog):
            ok_button.setEnabled(True)
            if(self.exported_thread.running):
                affichage = MessagePopUp(self) 
                dialog.accept()

        ok_button.clicked.connect(on_ok)
        cancel_button.clicked.connect(on_cancel)

        dialog.setLayout(dialog_layout)
        dialog.exec()

    def get_images(self):
        stock_images = []

        # Milieu de chaque scène
        stock_frames = [btn_data["frame1"] + (btn_data["frame2"] - btn_data["frame1"]) // 2 for btn_data in self.seg.display.stock_button]

        cap = cv2.VideoCapture(self.vlc.path_of_media)
        if not cap.isOpened():
            print("Impossible d'ouvrir la vidéo.")
            return

        frame_idx = 0
        stock_frames_set = set(stock_frames)

        while True:
            ret, frame = cap.read()
            if not ret:
                break
            if frame_idx in stock_frames_set:
                stock_images.append(frame)
            frame_idx += 1
        cap.release()

        return stock_images

    def size_of_img(self,img):
        max_width = 300
        max_height = 200

        _, img_bytes = cv2.imencode('.png', img)
        img_stream = BytesIO(img_bytes.tobytes())
        height, width, _ = img.shape
        width_ratio = max_width / float(width)
        height_ratio = max_height / float(height)
        ratio = min(width_ratio, height_ratio, 1.0)
        new_width = int(width * ratio)
        new_height = int(height * ratio)
        return img_stream,new_width,new_height
    
    def export_tagImages(self, callback=None):
        stock_images = self.get_images()
        try:
            tagImagesdir_path = os.path.join(self.file_path, "TagImages")
            print(f"Path du dossier : {tagImagesdir_path}")
            if os.path.exists(tagImagesdir_path):
                shutil.rmtree(tagImagesdir_path)
            os.makedirs(tagImagesdir_path)

            for idx, btn_data in enumerate(self.seg.display.stock_button):
                # time_str = time code du milieu de la scène 
                time_str = self.time_manager.frame_to_m(btn_data["frame1"] + (btn_data["frame2"] - btn_data["frame1"]) // 2)
                time_str = self.time_manager.timecodename(time_str)
                tagImage_path = os.path.join(tagImagesdir_path, f"TagImage{idx+1}_{time_str}.png")
                
                if idx < len(stock_images):
                    tagImage = stock_images[idx]

                    tagImage_stream, _, _ = self.size_of_img(tagImage)
                    tagImage_stream.seek(0)

                    with open(tagImage_path, "wb") as f:
                        f.write(tagImage_stream.read())
                    print(f"Image exportée : {tagImage_path}")
                else:
                    print(f"Aucune image disponible pour le timecode {time_str}")
        except Exception as e:
            print(f"Erreur lors de l'exportation des images : {e}")
    
    def export_txtV1(self, callback=None):
        self.file_path = os.path.join(self.file_path, f"{self.title}.txt")  # Correction de l'extension

        try:
            with open(self.file_path, "w", encoding="utf-8") as fichier:
                total_plans = len(self.seg.display.stock_button)
                fichier.write(f"{total_plans}\n")

                for btn_data in self.seg.display.stock_button:
                    if not callback():
                        print("Exportation annulée par l'utilisateur.")
                        return
                    button = btn_data["button"]
                    line = f"{button.text()}"
                    fichier.write(line + "\n")

            print(f"Fichier TXT enregistré : {self.file_path}")
        except Exception as e:
            print(f"Erreur lors de l'exportation TXT : {e}")

    def export_txt(self, callback=None):
        self.file_path = os.path.join(self.file_path, f"{self.title}.txt")

        try:
            with open(self.file_path, "w", encoding="utf-8") as fichier:
                fichier.write("=== Étude cinématographique ===\n\n")

                total_plans = len(self.seg.display.stock_button)
                fichier.write(f"Nombre total de plans : {total_plans}\n\n")

                for idx, btn_data in enumerate(self.seg.display.stock_button):
                    if not callback():
                        print("Exportation annulée par l'utilisateur.")
                        return

                    button = btn_data["button"]
                    time_str = self.time_manager.m_to_hmsf(btn_data["time"])
                    end_str = self.time_manager.m_to_hmsf(btn_data["end"] - btn_data["time"])

                    # En-tête du plan
                    fichier.write(f"- [Plan {idx+1}] {button.text()} -> Début : {time_str} / Durée : {end_str}\n")

                    # Notes associées (seulement si présentes)
                    for note_widget in self.seg.display.button_notes.get(button, []):
                        note_text = note_widget.toPlainText().strip()
                        if note_text:
                            fichier.write(f"{note_text}\n")

                    # Ligne vide entre les plans
                    fichier.write("\n")

                print(f"Fichier TXT enregistré : {self.file_path}")

        except Exception as e:
            print(f"Erreur lors de l'exportation TXT : {e}")


    def export_pdf(self,callback=None):
        self.file_path = os.path.join(self.file_path, f"{self.title}.pdf")

        stock_images=self.get_images()

        try:
            doc = SimpleDocTemplate(self.file_path, pagesize=A4)
            elements = []
            elements.append(Paragraph("Étude cinématographique", self.title_style))

            for idx, btn_data in enumerate(self.seg.display.stock_button):
                if not callback():
                    print("Exportation annulée par l'utilisateur.")
                    return
                
                button = btn_data["button"]
                time_str = self.time_manager.m_to_hmsf(btn_data["time"])
                end_str = self.time_manager.m_to_hmsf(btn_data["end"] - btn_data["time"])
                elements.append(Paragraph(f"- [Plan {idx+1}] {button.text()} -> Début : {time_str} / Durée : {end_str}", self.subtitle_style))
                
                for note_widget in self.seg.display.button_notes.get(button, []):
                    note_text = note_widget.toPlainText()
                    self.put_multiline_text(elements, note_text)
                
                if idx < len(stock_images):
                    img = stock_images[idx]
                    img_stream,new_width,new_height=self.size_of_img(img)
                    img_obj = PDFImage(img_stream, width=new_width, height=new_height)
                    elements.append(img_obj)
            
            doc.build(elements)
            print(f"Fichier PDF enregistré : {self.file_path}")
        except Exception as e:
            print(f"Erreur lors de l'exportation PDF : {e}")

    def export_pptx(self, callback=None):
        if callback is None:
            def _cb(): return True
            callback = _cb

        self.file_path = os.path.join(self.file_path, f"{self.title}.pptx")
        logging.info(f"[export_pptx] Création du pptx : {self.file_path}")
        stock_images = self.get_images()

        try:
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            blank_slide_layout = prs.slide_layouts[6]
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            slide = prs.slides.add_slide(blank_slide_layout)

            textbox_width = slide_width - Inches(1)
            textbox_left = (slide_width - textbox_width)/2

            title_height = Inches(2)
            title_box = slide.shapes.add_textbox(textbox_left, (slide_height - title_height)/2, textbox_width, title_height)
            title_tf = title_box.text_frame
            title_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            title_p = title_tf.paragraphs[0]
            title_p.alignment = PP_ALIGN.CENTER
            title_p.font.bold = True
            title_p.font.size = Pt(32)
            title_p.font.color.rgb = PPTXColor.from_string("EB4034")
            title_p.text = "Étude cinématographique"

            plan_box_top = Inches(0.5)
            plan_box_height = Inches(0.5)
            note_box_height = Inches(1)
            offset_between_boxes = Inches(0.2)

            for idx, btn_data in enumerate(self.seg.display.stock_button):
                if not callback():
                    print("Exportation annulée par l'utilisateur.")
                    return
                logging.info(f"[export_pptx] Création slide {idx}")

                button = btn_data["button"]
                time_str = self.time_manager.m_to_hmsf(btn_data["time"])
                end_str = self.time_manager.m_to_hmsf(btn_data["end"] - btn_data["time"])
                slide = prs.slides.add_slide(blank_slide_layout)

                y_offset = plan_box_top + plan_box_height + offset_between_boxes

                plan_box = slide.shapes.add_textbox(textbox_left, plan_box_top, textbox_width, plan_box_height)
                plan_tf = plan_box.text_frame
                plan_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

                plan_p = plan_tf.paragraphs[0]
                plan_p.font.bold = True
                plan_p.font.color.rgb = PPTXColor.from_string("1F497D")
                plan_p.text = f"[Plan {idx+1}] {button.text()} - Début : {time_str} / Durée : {end_str}"

                if idx < len(stock_images):
                    img = stock_images[idx]
                    img_stream, _, _ = self.size_of_img(img)
                    if isinstance(img_stream, bytes):
                        img_stream = BytesIO(img_stream)
                    else:
                        img_stream.seek(0)
                    try:
                        logging.info(f"[export_pptx] Taille Slide : {slide_width} * {slide_height} (w*h)")

                        image = Image.open(img_stream)
                        img_w, img_h = image.size
                        logging.info(f"[export_pptx] Taille image : {img_w} * {img_h} (w*h)")

                        dpi = image.info.get('dpi', (96, 96))[0]
                        img_width_emu, img_height_emu = Inches(img_w / dpi), Inches(img_h / dpi)

                        max_w, max_h = slide_width - Inches(1), slide_height - plan_box_height - note_box_height - y_offset - offset_between_boxes
                        logging.info(f"[export_pptx] Taille max : {max_w} * {max_h} (w*h)")

                        scale = min(max_w / img_width_emu, max_h / img_height_emu, 1.0)
                        logging.info(f"[export_pptx] Scale : {scale}")

                        new_w, new_h = img_width_emu * scale, img_height_emu * scale
                        logging.info(f"[export_pptx] Taille scalée : {new_w} * {new_h} (w*h)")

                        image_io = BytesIO()
                        image.save(image_io, format="PNG")
                        image_io.seek(0)

                        x_pos, y_pos = (slide_width - new_w) / 2, y_offset
                        logging.info(f"[export_pptx] position x et y : {x_pos} * {y_pos}")

                        slide.shapes.add_picture(image_io, x_pos, y_pos, width=new_w, height=new_h)
                        logging.info(f"[export_pptx] Ajout image {idx}")

                        y_offset += new_h + offset_between_boxes
                        
                    except Exception as e:
                        print(f"Erreur lors du chargement de l'image : {e}")

                notes = self.seg.display.button_notes.get(button, [])
                for note_widget in notes:
                    note_box = slide.shapes.add_textbox(textbox_left, y_offset, textbox_width, note_box_height)
                    note_box.height = note_box_height
                    note_tf = note_box.text_frame
                    note_tf.word_wrap = True

                    note_p = note_tf.paragraphs[0]
                    note_p.alignment = PP_ALIGN.LEFT
                    note_p.font.size = Pt(16)

                    note_text = note_widget.toPlainText().strip()
                    if note_text:
                        note_p.text = note_text
                    else:
                        note_tf.auto_size = MSO_AUTO_SIZE.NONE
                        note_p.text = "\u00A0"
                        note_box.height = note_box_height

            prs.save(self.file_path)
            print(f"Fichier PPTX enregistré : {self.file_path}")
            logging.info(f"Fichier PPTX enregistré : {self.file_path}")

        except Exception as e:
            print(f"Erreur lors de l'exportation PPTX : {e}")

    def export_docx(self, callback=None):
        self.file_path = os.path.join(self.file_path, f"{self.title}.docx")
        try:
            doc = Document()
            titre_paragraphe = doc.add_paragraph()
            titre_paragraphe.alignment = WD_ALIGN_PARAGRAPH.CENTER
            titre_run = titre_paragraphe.add_run("Étude cinématographique")
            titre_run.font.size = Pt(24)
            titre_run.font.color.rgb = DOCXColor(0xFF, 0x00, 0x00)
            
            stock_images = self.get_images()
        
            for idx, btn_data in enumerate(self.seg.display.stock_button):
                if not callback():
                    print("Exportation annulée par l'utilisateur.")
                    return
                button = btn_data["button"]
                time_str = self.time_manager.m_to_hmsf(btn_data["time"])
                end_str = self.time_manager.m_to_hmsf(btn_data["end"] - btn_data["time"])
                doc.add_heading(f"- [Plan {idx+1}] {button.text()} -> Début : {time_str} / Durée : {end_str}", level=2)
                
                for note_widget in self.seg.display.button_notes.get(button, []):
                    note_text = note_widget.toPlainText()
                    doc.add_paragraph(note_text)
                    
                if idx < len(stock_images):
                    img = stock_images[idx]
                    img_stream,new_width,new_height=self.size_of_img(img)
                    img_docx_width = Inches(new_width / 96.0)
                    p = doc.add_paragraph()
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    p.add_run().add_picture(img_stream, width=img_docx_width)
            
            doc.save(self.file_path)
            print(f"Fichier DOCX enregistré : {self.file_path}")
        except Exception as e:
            print(f"Erreur lors de l'exportation DOCX : {e}")

    def export_odt(self, callback=None):
        self.file_path = os.path.join(self.file_path, f"{self.title}.odt")
        try:
            doc = OpenDocumentText()

            title_style = Style(name="Titre", family="paragraph")
            title_style.addElement(ParagraphProperties(textalign="center"))
            title_style.addElement(TextProperties(color="#FF0000", fontsize="24pt"))
            doc.styles.addElement(title_style)

            titre_paragraphe = P(stylename=title_style, text="Étude cinématographique")
            doc.text.addElement(titre_paragraphe)

            first_button_style = Style(name="PremierBouton", family="paragraph")
            first_button_style.addElement(TextProperties(color="#0000FF", fontsize="14pt"))
            doc.styles.addElement(first_button_style)

            for idx, btn_data in enumerate(self.seg.display.stock_button):
                if not callback():
                    print("Exportation annulée par l'utilisateur.")
                    return
                
                button = btn_data["button"]
                time_str = self.time_manager.m_to_hmsf(btn_data["time"])
                end_str = self.time_manager.m_to_hmsf(btn_data["end"] - btn_data["time"])
                doc.text.addElement(P(stylename=first_button_style,text=f"- {button.text()} -> Début : {time_str} / Durée : {end_str}"))
                
                for note_widget in self.seg.display.button_notes.get(button, []):
                    note_text = note_widget.toPlainText()
                    doc.text.addElement(P(text=note_text))

                doc.text.addElement(P())

            doc.save(self.file_path)
            print(f"Fichier ODT enregistré : {self.file_path}")
        except Exception as e:
            print(f"Erreur lors de l'exportation ODT : {e}")

    def put_multiline_text(self, elements, text):
        lines = text.split("\n")
        for line in lines:
            line = line.replace("\t", "&nbsp;&nbsp;&nbsp;&nbsp;")
            elements.append(Paragraph(line, self.note_style))

    def export_video(self, callback=None):
        try:
            #self.default_file_path = os.path.join(self.file_path, f"{self.title}.mp4")

            if self.chosen_file_path:
                if not self.chosen_file_path.lower().endswith(".mp4"):
                    self.chosen_file_path += ".mp4"
            self.project_manager.path_of_super = self.chosen_file_path

            temp_dir = tempfile.gettempdir()
            temp_video_path = os.path.join(temp_dir, "temp_video.mp4")

            # Dossier temporaire MoviePy
            moviepy_tempdir = os.path.join(temp_dir, "moviepy_temp")
            if not os.path.exists(moviepy_tempdir):
                os.makedirs(moviepy_tempdir, exist_ok=True)

            # Forcer les chemins pour FFMPEG et les fichiers temporaires
            mp_config.FFMPEG_BINARY = os.environ.get("IMAGEIO_FFMPEG_EXE", "ffmpeg")
            mp_config.TEMP_DIR = moviepy_tempdir
            os.chdir(moviepy_tempdir)

            cap = cv2.VideoCapture(self.vlc.path_of_media)
            if not cap.isOpened():
                print("Impossible d'ouvrir la vidéo.")
                return

            fourcc = cv2.VideoWriter_fourcc(*"mp4v")
            out = cv2.VideoWriter(temp_video_path, fourcc, 30.0, (int(cap.get(3)), int(cap.get(4))))

            cpt = 0
            while cap.isOpened():
                if not callback():
                    print("Exportation annulée par l'utilisateur.")
                    cap.release()
                    out.release()
                    if os.path.exists(temp_video_path):
                        os.remove(temp_video_path)
                    return

                ret, frame = cap.read()
                if not ret:
                    break

                active_texts = [
                    btn_data for btn_data in self.seg.display.stock_button
                    if btn_data["frame1"] <= cpt < btn_data["frame2"]
                ]
                for btn_data in active_texts:
                    button = btn_data["button"]
                    time_str = self.time_manager.m_to_hmsf(btn_data["time"])
                    end_str = self.time_manager.m_to_hmsf(btn_data["end"] - btn_data["time"])
                    txt = f"[Plan {self.seg.display.stock_button.index(btn_data)+1}] {button.text()}"
                    txt2 = f"Debut : {time_str} / Duree : {end_str}"
                    txt3 = [note_widget.toPlainText() for note_widget in self.seg.display.button_notes.get(button, [])]
                    height, width, _ = frame.shape
                    self.write_text_horizontal_on_video2(frame, txt, txt2, txt3, width)
                out.write(frame)
                cpt += 1

            cap.release()
            out.release()

            if not callback():
                print("Exportation annulée par l'utilisateur.")
                if os.path.exists(temp_video_path):
                    os.remove(temp_video_path)
                return

            os.environ["IMAGEIO_FFMPEG_EXE"] = imageio_ffmpeg.get_ffmpeg_exe()

            logging.info("[export_video] Récupération du clip vidéo pour l'image...")
            video_clip = VideoFileClip(temp_video_path)
            try:
                audio_clip = AudioFileClip(self.vlc.path_of_media)
                #audio_clip = VideoFileClip(self.vlc.path_of_media).audio
                logging.info("[export_video] Création du clip vidéo pour le son...")
                logging.info(f"[export_video] self.vlc.path_of_media : {self.vlc.path_of_media}")

                # Vérifications
                logging.info(f"Taille du clip vidéo : {video_clip.size}")
                logging.info(f"Durée du clip vidéo : {video_clip.duration}")
                logging.info(f"Durée de l'audio : {audio_clip.duration if audio_clip else 'aucun'}")

                #final_clip = video_clip.set_audio(audio_clip)
                final_clip = video_clip.with_audio(audio_clip)
            except Exception as e:
                print(f"⚠️ Impossible de charger l'audio : {e}")
                final_clip = video_clip
            
            final_clip.write_videofile(self.chosen_file_path, codec="libx264", audio_codec="aac", logger=ProgressBarLogger())
            logging.info(f"[export_video] Écriture faite")

            video_clip.close()
            audio_clip.close()
            final_clip.close()

            os.remove(temp_video_path)
            if os.path.exists(moviepy_tempdir):
                shutil.rmtree(moviepy_tempdir)

        except Exception as e:
            print(f"Erreur pendant l'export vidéo : {e}")

    def write_text_horizontal_on_video(self, frame, txt, txt2, txt3, max_width, line_spacing=24):
        final_txt = txt + " - " + txt2 + "\n"
        for text in txt3:
            lines = text.split("\n")
            final_txt += " - " + " / ".join(line.replace("\t", "   ") for line in lines)

        font_scale = 0.7
        thickness_outline = 4
        thickness_text = 1
        font = cv2.FONT_HERSHEY_SIMPLEX

        char_size = cv2.getTextSize("A", font, font_scale, thickness_text)[0][0]
        max_chars_per_line = (max_width + 50) // char_size

        wrapped_text = textwrap.fill(final_txt, width=max_chars_per_line)
        x, y = 50, 30
        for line in wrapped_text.split("\n"):
            cv2.putText(frame, line, (x, y), font, font_scale, (0, 0, 0), thickness_outline, cv2.LINE_AA)
            cv2.putText(frame, line, (x, y), font, font_scale, (255, 255, 255), thickness_text, cv2.LINE_AA)
            y += line_spacing

    def write_text_horizontal_on_video2(self, frame, txt, txt2, txt3, max_width, line_spacing=24):
        font_scale = 0.7
        thickness_outline = 4
        thickness_text = 1
        font = cv2.FONT_HERSHEY_SIMPLEX

        char_size = cv2.getTextSize("A", font, font_scale, thickness_text)[0][0]
        max_chars_per_line = (max_width + 50) // char_size

        lines_to_draw = [txt + " - " + txt2]

        for text in txt3:
            for subline in text.split("\n"):
                processed_line = subline.replace("\t", "   ")
                lines_to_draw.append(processed_line)

        wrapped_lines = []
        for line in lines_to_draw:
            wrapped_lines.extend(textwrap.wrap(line, width=max_chars_per_line))

        x, y = 50, 30
        for line in wrapped_lines:
            cv2.putText(frame, line, (x, y), font, font_scale, (0, 0, 0), thickness_outline, cv2.LINE_AA)
            cv2.putText(frame, line, (x, y), font, font_scale, (255, 255, 255), thickness_text, cv2.LINE_AA)
            y += line_spacing


    def write_text_on_video(self, frame, txt, txt2, txt3, decalage):
        cv2.putText(frame, txt, (50, decalage + 50), cv2.FONT_HERSHEY_SIMPLEX, 1, (0, 0, 0), 5, cv2.LINE_AA)
        cv2.putText(frame, txt, (50, decalage + 50), cv2.FONT_HERSHEY_SIMPLEX, 1, (255, 255, 255), 2, cv2.LINE_AA)
        cv2.putText(frame, txt2, (50, decalage + 80), cv2.FONT_HERSHEY_SIMPLEX, 0.7, (0, 0, 0), 5, cv2.LINE_AA)
        cv2.putText(frame, txt2, (50, decalage + 80), cv2.FONT_HERSHEY_SIMPLEX, 0.7, (255, 255, 255), 2, cv2.LINE_AA)
        val = decalage + 80
        for i, text in enumerate(txt3):
            val = decalage + 110 + (i * 20)
            self.put_multiline_text_video(frame, text, (50, val))
        return val + 10

    def put_multiline_text_video(self, frame, text, position, font_scale=0.5, line_spacing=15):
        lines = text.split("\n")
        x, y = position
        for i, line in enumerate(lines):
            cv2.putText(frame, line.replace("\t", "    "), (x, y + i * line_spacing),
                        cv2.FONT_HERSHEY_SIMPLEX, font_scale, (0, 0, 0), 5, cv2.LINE_AA)
            cv2.putText(frame, line.replace("\t", "    "), (x, y + i * line_spacing),
                        cv2.FONT_HERSHEY_SIMPLEX, font_scale, (255, 255, 255), 2, cv2.LINE_AA)
